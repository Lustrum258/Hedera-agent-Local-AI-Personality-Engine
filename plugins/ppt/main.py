"""
PPT 制作插件
根据文字内容生成 .pptx 文件，保存到桌面。
支持：标题、正文、多级列表、多张幻灯片。
"""

import os
import re
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from hedera.plugin.base import PluginBase


# ─── 配色方案 ───

_THEMES = {
    "default": {
        "bg": RGBColor(0x1A, 0x1A, 0x2E),       # 深蓝黑
        "title": RGBColor(0xE0, 0xE0, 0xFF),     # 浅紫白
        "subtitle": RGBColor(0xA0, 0xC4, 0xFF),  # 浅蓝
        "text": RGBColor(0xCC, 0xCC, 0xDD),      # 浅灰紫
        "accent": RGBColor(0x7C, 0x83, 0xEB),    # 紫罗兰
        "highlight": RGBColor(0xFF, 0xD5, 0x4F), # 金色高亮
    },
    "light": {
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "title": RGBColor(0x1A, 0x1A, 0x2E),
        "subtitle": RGBColor(0x4A, 0x4A, 0x6A),
        "text": RGBColor(0x33, 0x33, 0x44),
        "accent": RGBColor(0x7C, 0x83, 0xEB),
        "highlight": RGBColor(0xE6, 0xA8, 0x17),
    },
    "dark": {
        "bg": RGBColor(0x0D, 0x0D, 0x0D),
        "title": RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle": RGBColor(0xBB, 0xBB, 0xBB),
        "text": RGBColor(0x99, 0x99, 0x99),
        "accent": RGBColor(0x66, 0x66, 0xFF),
        "highlight": RGBColor(0xFF, 0xCC, 0x00),
    },
}


def _make_ppt(title: str, slides_data: list[dict], output_path: str, theme: str = "default") -> str:
    """
    生成 PPTX 文件。
    slides_data: [
        {"title": "标题", "content": "正文\n- 要点1\n- 要点2", "type": "bullet|text|two_column"},
        ...
    ]
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    colors = _THEMES.get(theme, _THEMES["default"])

    blank_layout = prs.slide_layouts[6]  # Blank layout

    for idx, slide_info in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)

        # 背景色
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = colors["bg"]

        slide_title = slide_info.get("title", "")
        slide_content = slide_info.get("content", "")
        slide_type = slide_info.get("type", "bullet")

        # 标题
        if slide_title:
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(1.0))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_title
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = colors["title"]
            p.alignment = PP_ALIGN.LEFT

            # 标题下分隔线
            line = slide.shapes.add_shape(
                1,  # MSO_SHAPE.RECTANGLE
                Inches(0.8), Inches(1.3), Inches(11.7), Pt(3)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = colors["accent"]
            line.line.fill.background()

        # 内容
        if slide_content:
            content_top = Inches(1.8) if slide_title else Inches(0.8)
            txBox = slide.shapes.add_textbox(Inches(0.8), content_top, Inches(11.7), Inches(5.0))
            tf = txBox.text_frame
            tf.word_wrap = True

            lines = slide_content.strip().split("\n")
            first = True
            for line in lines:
                stripped = line.strip()
                if not stripped:
                    continue

                if first:
                    p = tf.paragraphs[0]
                    first = False
                else:
                    p = tf.add_paragraph()

                # 判断层级
                if stripped.startswith("## "):
                    p.text = stripped[3:]
                    p.font.size = Pt(24)
                    p.font.bold = True
                    p.font.color.rgb = colors["subtitle"]
                    p.space_before = Pt(16)
                    p.space_after = Pt(8)

                elif stripped.startswith("### "):
                    p.text = stripped[4:]
                    p.font.size = Pt(20)
                    p.font.bold = True
                    p.font.color.rgb = colors["accent"]
                    p.space_before = Pt(12)
                    p.space_after = Pt(4)

                elif stripped.startswith("- ") or stripped.startswith("* "):
                    p.text = stripped[2:]
                    p.font.size = Pt(16)
                    p.font.color.rgb = colors["text"]
                    # 缩进
                    p.level = 1
                    p.space_before = Pt(4)
                    p.space_after = Pt(2)

                elif stripped.startswith("  - ") or stripped.startswith("  * "):
                    p.text = stripped[4:]
                    p.font.size = Pt(14)
                    p.font.color.rgb = colors["text"]
                    p.level = 2
                    p.space_before = Pt(2)

                elif stripped.startswith("> "):
                    p.text = stripped[2:]
                    p.font.size = Pt(14)
                    p.font.italic = True
                    p.font.color.rgb = colors["highlight"]
                    p.space_before = Pt(6)
                    p.space_after = Pt(6)

                else:
                    p.text = stripped
                    p.font.size = Pt(18)
                    p.font.color.rgb = colors["text"]
                    p.space_before = Pt(6)
                    p.space_after = Pt(4)

                p.alignment = PP_ALIGN.LEFT

        # 页码
        page_num = slide.shapes.add_textbox(Inches(12.0), Inches(7.0), Inches(1.0), Inches(0.4))
        pf = page_num.text_frame
        p = pf.paragraphs[0]
        p.text = f"{idx + 1} / {len(slides_data)}"
        p.font.size = Pt(10)
        p.font.color.rgb = colors["accent"]
        p.alignment = PP_ALIGN.RIGHT

    prs.save(output_path)
    return output_path


def _parse_content_to_slides(content: str, default_title: str = "") -> list[dict]:
    """
    将纯文本解析为幻灯片数据结构。
    支持两种格式：
    1. `---` 分隔的多张幻灯片
    2. 每段标题自动成页
    """
    slides_data = []

    # 尝试按 `---` 或 `===` 分割
    blocks = re.split(r'\n---+\n|\n===+\n', content.strip())

    if len(blocks) <= 1 and not content.strip().startswith("---"):
        # 无分隔符：整段作为一页
        lines = content.strip().split("\n")
        title = default_title or lines[0] if lines else "Hedera PPT"
        body = "\n".join(lines[1:]) if len(lines) > 1 else ""
        slides_data.append({
            "title": title,
            "content": body if body else "（内容）",
            "type": "bullet",
        })
    else:
        for block in blocks:
            block = block.strip()
            if not block:
                continue
            lines = block.split("\n")
            # 第一行作为标题（除非以特殊前缀开头）
            first = lines[0].strip()
            if first.startswith("- ") or first.startswith("* ") or first.startswith("> "):
                # 第一行也是内容的一部分
                title = default_title or f"第 {len(slides_data) + 1} 页"
                body = "\n".join(lines)
            elif len(lines) >= 2 and lines[1].strip().startswith("="):
                # Markdown 标题格式
                title = first
                body = "\n".join(lines[2:])
            else:
                title = first
                body = "\n".join(lines[1:]) if len(lines) > 1 else ""

            slides_data.append({
                "title": title,
                "content": body,
                "type": "bullet",
            })

    if not slides_data:
        slides_data.append({
            "title": default_title or "Hedera PPT",
            "content": content,
            "type": "bullet",
        })

    return slides_data


def _generate_ppt(title: str, content: str, theme: str = "default", output_name: str = "") -> dict:
    """生成 PPT 文件"""
    try:
        slides_data = _parse_content_to_slides(content, title)

        if not output_name:
            safe_title = re.sub(r'[\\/:*?"<>|]', '_', title)[:30]
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_name = f"{safe_title}_{timestamp}.pptx"

        desktop = os.path.join(os.path.expanduser("~"), "Desktop")
        output_path = os.path.join(desktop, output_name)

        result_path = _make_ppt(title, slides_data, output_path, theme)

        return {
            "success": True,
            "path": result_path,
            "filename": output_name,
            "slides": len(slides_data),
            "message": f"已生成 {len(slides_data)} 页 PPT → {output_name}",
        }
    except Exception as e:
        return {"success": False, "error": str(e)}


# ─── 插件类 ───


class PPTPlugin(PluginBase):
    name = "ppt"
    description = "根据文字内容生成 PowerPoint 演示文稿"
    keywords = ["ppt", "幻灯片", "演示文稿", "pptx", "powerpoint", "讲稿", "生成pp"]
    commands = ["/ppt"]

    def process(self, message: str, context: dict = None) -> str | None:
        msg = message.strip()

        # /ppt 命令
        if msg.startswith("/ppt "):
            content = msg[5:].strip()
            if not content:
                return "格式：/ppt <标题>\n---\n<内容>\n\n用 --- 分隔多张幻灯片，- 表示列表项。"
            lines = content.split("\n")
            title = lines[0]
            body = "\n".join(lines[1:])
            result = _generate_ppt(title, body)
            if result["success"]:
                return f"✅ {result['message']}\n📁 {result['path']}"
            return f"❌ 生成失败: {result['error']}"

        # 关键词匹配 —— 检查是否描述了 PPT 需求
        has_ppt_kw = any(kw in msg for kw in ["ppt", "幻灯片", "演示文稿", "讲稿", "pptx"])
        if has_ppt_kw and len(msg) > 15:
            # 非命令模式：提取内容（交给 LLM 处理更合适）
            return None  # 让 LLM 调用工具处理

        return None

    def get_tools(self) -> list[dict]:
        return [{
            "name": "create_ppt",
            "description": "根据文字内容生成 PowerPoint 演示文稿（.pptx），可指定主题色：default(深色)/light(浅色)/dark(纯黑)",
            "fn": _generate_ppt,
            "parameters": {
                "type": "object",
                "properties": {
                    "title": {
                        "type": "string",
                        "description": "PPT 标题（第一页标题）",
                    },
                    "content": {
                        "type": "string",
                        "description": "幻灯片内容。用 `---` 分隔多张幻灯片。\n"
                                       "- 每段第一行是该页标题\n"
                                       "- `- 文本` 表示列表项\n"
                                       "- `## 文本` 表示二级标题\n"
                                       "- `> 文本` 表示引用/高亮\n"
                                       "- 示例：\n"
                                       "用户画像分析\n---\n## 核心用户群\n- 年龄 25-40\n- 一线城市\n---\n## 用户痛点\n- 时间碎片化\n- 信息过载",
                    },
                    "theme": {
                        "type": "string",
                        "enum": ["default", "light", "dark"],
                        "description": "主题色：default(深蓝黑), light(白底), dark(纯黑)",
                        "default": "default",
                    },
                    "output_name": {
                        "type": "string",
                        "description": "可选：自定义文件名（不含路径，保存到桌面）",
                    },
                },
                "required": ["title", "content"],
            },
        }]

    def get_system_prompt_modifier(self) -> str:
        return (
            "【工具】`create_ppt(title, content, theme, output_name)` — 根据文字生成 PPT。\n"
            "  content 中用 `---` 分隔多页，`- ` 表示列表，`## ` 表示小标题。\n"
            "  用户要你「做个PPT」或「写成幻灯片」时，先设计内容再调工具生成，直接问内容够不够详细。"
        )
